import gradio as gr # user interface
import os
import glob
import logging
from dotenv import load_dotenv
from pptx import Presentation
import PyPDF2  # Add PDF library
import re      # For text file processing
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
from langchain_google_genai import ChatGoogleGenerativeAI
import google.generativeai as genai

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("chatbot.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ctse_chatbot")

# Load environment variables from .env file
try:
    load_dotenv()
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
    if not GOOGLE_API_KEY:
        raise ValueError("GOOGLE_API_KEY not found in environment variables")
    
    # Configure Gemini API
    genai.configure(api_key=GOOGLE_API_KEY)
    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
    
    # Get other environment variables with defaults
    VECTOR_STORE_DIR = os.getenv("VECTOR_STORE_DIR", "./vector_store")
    LECTURES_DIR = os.getenv("LECTURES_DIR", "./lectures")
    MODEL_NAME = os.getenv("MODEL_NAME", "gemini-2.0-flash")
    TEMPERATURE = float(os.getenv("TEMPERATURE", "0.7"))
    TOP_P = float(os.getenv("TOP_P", "0.85"))
    MAX_TOKENS = int(os.getenv("MAX_TOKENS", "1024"))
    EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
    
    logger.info("Environment variables loaded successfully")
except Exception as e:
    logger.error(f"Error loading environment variables: {str(e)}")
    raise

# Function to extract text from PowerPoint files
def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint file with error handling."""
    try:
        prs = Presentation(pptx_path)
        text_content = []
        
        # Extract slide number for context
        slide_number = 1
        
        for slide in prs.slides:
            try:
                slide_text = f"Slide {slide_number}: "
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                text_content.append(slide_text.strip())
                slide_number += 1
            except Exception as slide_error:
                logger.warning(f"Error processing slide {slide_number} in {pptx_path}: {str(slide_error)}")
                text_content.append(f"Slide {slide_number}: [Error extracting content]")
                slide_number += 1
        
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to extract text from {pptx_path}: {str(e)}")
        return f"Error processing file {os.path.basename(pptx_path)}: {str(e)}"

# New function to extract text from PDF files
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file with error handling."""
    try:
        text_content = []
        
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            page_count = len(pdf_reader.pages)
            
            for page_num in range(page_count):
                try:
                    page = pdf_reader.pages[page_num]
                    page_text = f"Page {page_num + 1}: {page.extract_text()}"
                    text_content.append(page_text.strip())
                except Exception as page_error:
                    logger.warning(f"Error processing page {page_num + 1} in {pdf_path}: {str(page_error)}")
                    text_content.append(f"Page {page_num + 1}: [Error extracting content]")
        
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to extract text from {pdf_path}: {str(e)}")
        return f"Error processing file {os.path.basename(pdf_path)}: {str(e)}"

# New function to extract text from plain text files
def extract_text_from_txt(txt_path):
    """Extract text from a plain text file with error handling."""
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            content = file.read()
            
        # Optional: Add line numbering or structure to the content
        lines = content.split('\n')
        text_content = []
        
        for i, line in enumerate(lines):
            if line.strip():  # Skip empty lines
                text_content.append(f"Line {i+1}: {line}")
        
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to extract text from {txt_path}: {str(e)}")
        return f"Error processing file {os.path.basename(txt_path)}: {str(e)}"

# Enhanced function to load document files (PPTX, PDF, TXT)
def load_document_files(directory_path):
    """Load all document files (PPTX, PDF, TXT) from a directory and extract their text."""
    try:
        # Check if directory exists
        if not os.path.exists(directory_path):
            logger.error(f"Directory not found: {directory_path}")
            raise FileNotFoundError(f"Directory not found: {directory_path}")
        
        # Find all files of the specified types
        pptx_files = glob.glob(os.path.join(directory_path, "*.pptx"))
        pdf_files = glob.glob(os.path.join(directory_path, "*.pdf"))
        txt_files = glob.glob(os.path.join(directory_path, "*.txt"))
        
        all_files = pptx_files + pdf_files + txt_files
        
        if not all_files:
            logger.warning(f"No document files found in {directory_path}")
        
        all_text = []
        file_sources = []
        processed_count = 0
        
        # Process PowerPoint files
        for file_path in pptx_files:
            try:
                file_name = os.path.basename(file_path)
                logger.info(f"Processing PPTX: {file_name}")
                text = extract_text_from_pptx(file_path)
                # Add file source for better context
                text = f"Source: {file_name}\n\n{text}"
                all_text.append(text)
                file_sources.append(file_name)
                processed_count += 1
            except Exception as file_error:
                logger.error(f"Error processing file {file_path}: {str(file_error)}")
        
        # Process PDF files
        for file_path in pdf_files:
            try:
                file_name = os.path.basename(file_path)
                logger.info(f"Processing PDF: {file_name}")
                text = extract_text_from_pdf(file_path)
                # Add file source for better context
                text = f"Source: {file_name}\n\n{text}"
                all_text.append(text)
                file_sources.append(file_name)
                processed_count += 1
            except Exception as file_error:
                logger.error(f"Error processing file {file_path}: {str(file_error)}")
        
        # Process TXT files
        for file_path in txt_files:
            try:
                file_name = os.path.basename(file_path)
                logger.info(f"Processing TXT: {file_name}")
                text = extract_text_from_txt(file_path)
                # Add file source for better context
                text = f"Source: {file_name}\n\n{text}"
                all_text.append(text)
                file_sources.append(file_name)
                processed_count += 1
            except Exception as file_error:
                logger.error(f"Error processing file {file_path}: {str(file_error)}")
        
        logger.info(f"Processed {processed_count} out of {len(all_files)} document files")
        return all_text, file_sources
    
    except Exception as e:
        logger.error(f"Failed to load document files: {str(e)}")
        raise

# Process and split text
def process_text(texts):
    """Split the text into smaller chunks for better retrieval with error handling."""
    try:
        if not texts:
            logger.warning("No texts provided for processing")
            return []
            
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        
        all_splits = []
        for i, text in enumerate(texts):
            try:
                splits = text_splitter.split_text(text)
                all_splits.extend(splits)
            except Exception as text_error:
                logger.warning(f"Error splitting text chunk {i}: {str(text_error)}")
        
        logger.info(f"Created {len(all_splits)} text chunks from {len(texts)} documents")
        return all_splits
    
    except Exception as e:
        logger.error(f"Failed to process texts: {str(e)}")
        raise

# Create Vector Store
def create_vector_store(text_chunks):
    """Create a vector store using FAISS for efficient similarity search with error handling."""
    try:
        if not text_chunks:
            logger.warning("No text chunks provided for vector store creation")
            raise ValueError("No text chunks provided for vector store creation")
            
        # Create directory for vector store if it doesn't exist
        os.makedirs(VECTOR_STORE_DIR, exist_ok=True)
        
        logger.info(f"Initializing embeddings model: {EMBEDDING_MODEL}")
        embeddings = HuggingFaceEmbeddings(
            model_name=EMBEDDING_MODEL,
            model_kwargs={'device': 'cpu'}
        )
        
        logger.info(f"Creating vector store with {len(text_chunks)} chunks")
        vector_store = FAISS.from_texts(text_chunks, embeddings)
        
        # Save the vector store for future use
        vector_store_path = os.path.join(VECTOR_STORE_DIR, "faiss_index")
        vector_store.save_local(vector_store_path)
        logger.info(f"Vector store saved to {vector_store_path}")
        
        return vector_store
    
    except Exception as e:
        logger.error(f"Failed to create vector store: {str(e)}")
        raise

# Try to load existing vector store or create a new one
def get_vector_store(text_chunks=None):
    """Get vector store - either load existing one or create new one."""
    try:
        vector_store_path = os.path.join(VECTOR_STORE_DIR, "faiss_index")
        
        if os.path.exists(vector_store_path) and not text_chunks:
            try:
                logger.info("Loading existing vector store")
                embeddings = HuggingFaceEmbeddings(
                    model_name=EMBEDDING_MODEL,
                    model_kwargs={'device': 'cpu'}
                )
                vector_store = FAISS.load_local(vector_store_path, embeddings)
                logger.info("Vector store loaded successfully")
                return vector_store
            except Exception as load_error:
                logger.warning(f"Failed to load vector store: {str(load_error)}. Will create new one.")
        
        if text_chunks:
            return create_vector_store(text_chunks)
        else:
            raise ValueError("No text chunks provided and no existing vector store found")
    
    except Exception as e:
        logger.error(f"Error with vector store: {str(e)}")
        raise

# Set Up Gemini LLM
def setup_gemini_llm():
    """Set up the Gemini Pro LLM via API with error handling."""
    try:
        logger.info(f"Setting up Gemini LLM with model: {MODEL_NAME}")
        llm = ChatGoogleGenerativeAI(
            model=MODEL_NAME,
            temperature=TEMPERATURE,
            top_p=TOP_P,
            max_output_tokens=MAX_TOKENS,
            convert_system_message_to_human=True
        )
        
        return llm
    
    except Exception as e:
        logger.error(f"Failed to set up Gemini LLM: {str(e)}")
        raise

# Create Conversational Retrieval Chain
def create_chatbot(vector_store, llm, memory):
    """Create the conversational retrieval chain with error handling."""
    try:        
        qa_chain = ConversationalRetrievalChain.from_llm(
            llm=llm,
            retriever=vector_store.as_retriever(search_kwargs={"k": 3}),
            memory=memory,
            return_source_documents=True
        )
        
        logger.info("Conversational retrieval chain created successfully")
        return qa_chain
    
    except Exception as e:
        logger.error(f"Failed to create conversational chain: {str(e)}")
        raise

# Update the initialize_chatbot function to use the new load_document_files function
def initialize_chatbot(doc_directory=None):
    """Initialize the chatbot with proper error handling."""
    try:
        # Use environment variable if no directory provided
        if doc_directory is None:
            doc_directory = LECTURES_DIR
        
        # Ensure the lectures directory exists
        if not os.path.exists(doc_directory):
            os.makedirs(doc_directory, exist_ok=True)
            logger.warning(f"Created lectures directory: {doc_directory}")
        
        # Check if we have an existing vector store
        vector_store_path = os.path.join(VECTOR_STORE_DIR, "faiss_index")
        
        if os.path.exists(vector_store_path):
            try:
                # Try to load existing vector store
                vector_store = get_vector_store()
                logger.info("Using existing vector store")
                
                # Just get the list of files
                _, file_sources = load_document_files(doc_directory)
            except Exception as vs_error:
                logger.warning(f"Error loading vector store: {str(vs_error)}. Creating new one.")
                # Process files and create new vector store
                all_documents, file_sources = load_document_files(doc_directory)
                logger.info(f"Loaded {len(all_documents)} document files")
                
                # Process and split the text
                text_chunks = process_text(all_documents)
                logger.info(f"Created {len(text_chunks)} text chunks")
                
                # Create vector store
                vector_store = create_vector_store(text_chunks)
        else:
            # No existing vector store, so create one
            all_documents, file_sources = load_document_files(doc_directory)
            logger.info(f"Loaded {len(all_documents)} document files")
            
            # Process and split the text
            text_chunks = process_text(all_documents)
            logger.info(f"Created {len(text_chunks)} text chunks")
            
            # Create vector store
            vector_store = create_vector_store(text_chunks)
        
        # Setup memory to maintain conversation history
        memory = ConversationBufferMemory(
            memory_key="chat_history",
            output_key="answer",
            return_messages=True
        )
        
        # Set up LLM and QA chain
        llm = setup_gemini_llm()
        qa_chain = create_chatbot(vector_store, llm, memory)
        logger.info("Chatbot initialized successfully")
        
        return qa_chain, file_sources
    
    except Exception as e:
        logger.error(f"Failed to initialize chatbot: {str(e)}")
        raise


# Main chatbot response function for TUPLE format (the original format)
def get_response(message, history):
    """Get response from the chatbot with error handling."""
    if not message.strip():
        history = history or []
        history.append((message, "Please enter a question."))
        return history
        
    try:
        # Convert message to the format expected by the chain
        logger.info(f"Processing question: {message}")
        result = qa_chain({"question": message})
        answer = result["answer"]
        
        # Format source information
        sources = []
        for doc in result["source_documents"]:
            try:
                source = doc.metadata.get("source", "Unknown source")
                if source not in sources and "Source:" in doc.page_content:
                    source_line = [line for line in doc.page_content.split('\n') if 'Source:' in line]
                    if source_line:
                        source = source_line[0].replace('Source:', '').strip()
                
                if source not in sources:
                    sources.append(source)
            except Exception as doc_error:
                logger.warning(f"Error processing document source: {str(doc_error)}")
        
        # Add sources to the response
        if sources:
            source_text = "\n\n**Sources:**\n"
            for src in sources:
                if src and src != "Unknown source":
                    source_text += f"- {src}\n"
            
            if source_text != "\n\n**Sources:**\n":
                answer += source_text
        
        logger.info("Successfully generated response")
        
        # Update history with the new message pair
        history = history or []
        history.append((message, answer))
        return history
    
    except Exception as e:
        logger.error(f"Error generating response: {str(e)}")
        error_message = f"Error: {str(e)}\nPlease try again with a different question."
        history = history or []
        history.append((message, error_message))
        return history

# Error response for the TUPLE format
def error_response(message, history):
    history = history or []
    history.append((message, "Chatbot initialization failed. Please check the configuration and restart the application."))
    return history

# Try to initialize the chatbot - wrapped in try/except for robust error handling
try:
    # Initialize the chatbot
    qa_chain, available_files = initialize_chatbot()
    logger.info("Chatbot initialization complete")
except Exception as init_error:
    logger.critical(f"Critical error initializing chatbot: {str(init_error)}")
    # Set these to None so we can handle in the UI
    qa_chain = None
    available_files = []

# Define Gradio interface
with gr.Blocks(css="footer {visibility: hidden}") as demo:
    gr.Markdown(
        """
        # CTSE Lecture Notes Chatbot
        
        Welcome to the **Current Trends in Software Engineering** Lecture Notes Chatbot! 
        Ask any questions about the lecture content, and I'll try to answer based on the available lecture notes.

        Supported file formats: PowerPoint (PPTX), PDF, and Text (TXT) files.
        """
    )
    
    # Show error message if initialization failed
    if qa_chain is None:
        gr.Markdown(
            """
            ⚠️ **Error: Chatbot initialization failed.**
            
            Please check the log file for details and make sure:
            1. Your .env file is properly configured
            2. The Google API key is valid
            3. Lecture files are accessible
            
            The chatbot interface is still available but may not function correctly.
            """
        )
    
    with gr.Row():
        with gr.Column(scale=4):
            # IMPORTANT: Using standard chatbot without type="messages"
            chatbot = gr.Chatbot(height=500)
            with gr.Row():
                msg = gr.Textbox(
                    placeholder="Ask a question about CTSE lectures...",
                    show_label=False,
                    scale=9
                )
                submit = gr.Button("Send", scale=1)
            clear = gr.Button("Clear Chat")
        
        with gr.Column(scale=1):
            gr.Markdown("### Available Lecture Notes")
            file_list = gr.Dataframe(
                headers=["Lecture Files"],
                datatype=["str"],
                value=[[file] for file in available_files]
            )
    
    # Set up event handlers - only if chatbot was initialized successfully
    if qa_chain is not None:
        msg.submit(
            get_response,
            [msg, chatbot],
            [chatbot],
            queue=False
        ).then(
            lambda: "",
            None,
            [msg],
            queue=False
        )
        
        submit.click(
            get_response,
            [msg, chatbot],
            [chatbot],
            queue=False
        ).then(
            lambda: "",
            None,
            [msg],
            queue=False
        )
    else:
        # Display an error message when attempted to use
        msg.submit(
            error_response,
            [msg, chatbot],
            [chatbot],
            queue=False
        ).then(
            lambda: "",
            None,
            [msg],
            queue=False
        )
        
        submit.click(
            error_response,
            [msg, chatbot],
            [chatbot],
            queue=False
        ).then(
            lambda: "",
            None,
            [msg],
            queue=False
        )
    
    clear.click(lambda: [], None, chatbot, queue=False)

    gr.Markdown(
        """
        ### About This Chatbot
        This chatbot is powered by Google's Gemini AI and uses vector embeddings to search through CTSE lecture slides.
        It can help answer questions about cloud computing, microservices, DevOps, and other topics covered in the course.
        
        **Note**: The chatbot's knowledge is limited to the content of the lecture slides.
        """
    )

# Launch the Gradio interface
if __name__ == "__main__":
    try:
        logger.info("Starting Gradio interface")
        demo.launch(share=True)
    except Exception as launch_error:
        logger.critical(f"Failed to launch Gradio interface: {str(launch_error)}")